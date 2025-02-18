from flask import Flask, render_template, request, send_file, abort
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

app = Flask(__name__)
1
@app.route('/')
def index():
    return render_template('index.html')

@app.route('/generate', methods=['POST'])
def generate_ppt():
    try:
        components = {
            "Key Partners": request.form.get('key_partners'),
            "Key Activities": request.form.get('key_activities'),
            "Key Resources": request.form.get('key_resources'),
            "Value Propositions": request.form.get('value_propositions'),
            "Customer Relationships": request.form.get('customer_relationships'),
            "Channels": request.form.get('channels'),
            "Customer Segments": request.form.get('customer_segments'),
            "Cost Structure": request.form.get('cost_structure'),
            "Revenue Streams": request.form.get('revenue_streams')
        }

        # Create a presentation object
        prs = Presentation()

        # Add a title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Business Model Canvas"
        subtitle.text = "Generated by Flask and python-pptx"
        
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 128, 128)  # Teal color

        subtitle.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 100, 0)  # Dark green color

        # Add a slide for the Business Model Canvas structure
        slide_layout = prs.slide_layouts[5]  # Use a blank slide layout
        slide = prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
        
        # Define positions and dimensions for each block in the Business Model Canvas layout
        positions = {
            "Key Partners": (Inches(0.5), Inches(1), Inches(2.5), Inches(1.5)),
            "Key Activities": (Inches(3.2), Inches(1), Inches(2.5), Inches(1.5)),
            "Key Resources": (Inches(5.9), Inches(1), Inches(2.5), Inches(1.5)),
            "Value Propositions": (Inches(0.5), Inches(2.75), Inches(8), Inches(1.5)),
            "Customer Relationships": (Inches(0.5), Inches(4.5), Inches(2.5), Inches(1.5)),
            "Channels": (Inches(3.2), Inches(4.5), Inches(2.5), Inches(1.5)),
            "Customer Segments": (Inches(5.9), Inches(4.5), Inches(2.5), Inches(1.5)),
            "Cost Structure": (Inches(0.5), Inches(6.25), Inches(3.5), Inches(1.5)),
            "Revenue Streams": (Inches(5), Inches(6.25), Inches(3.5), Inches(1.5))
        }

        # Define colors for each block
        colors = {
            "Key Partners": RGBColor(255, 204, 204),  # Light Red
            "Key Activities": RGBColor(204, 255, 204),  # Light Green
            "Key Resources": RGBColor(204, 204, 255),  # Light Blue
            "Value Propositions": RGBColor(255, 255, 204),  # Light Yellow
            "Customer Relationships": RGBColor(255, 204, 255),  # Light Pink
            "Channels": RGBColor(204, 255, 255),  # Light Cyan
            "Customer Segments": RGBColor(255, 229, 204),  # Light Orange
            "Cost Structure": RGBColor(229, 204, 255),  # Light Purple
            "Revenue Streams": RGBColor(204, 229, 255)  # Light Sky Blue
        }

        # Add text boxes for each component with different colors
        for key, pos in positions.items():
            left, top, width, height = pos
            shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = colors[key]
            shape.line.color.rgb = RGBColor(0, 0, 0)  # Black border
            text_frame = shape.text_frame
            text_frame.word_wrap = True
            p = text_frame.add_paragraph()
            p.text = f"{key}:\n{components[key] if components[key] else 'No data provided'}"
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(0, 0, 0)  # Black color
            p.line_spacing = 1.5

        # Save the presentation
        current_dir = os.getcwd()
        ppt_path = os.path.join(current_dir, "business_model_canvas.pptx")
        prs.save(ppt_path)
        print(f"Presentation saved to {ppt_path}")  # Debugging

        # Check if the file exists
        if not os.path.exists(ppt_path):
            print("File not found after saving:", ppt_path)  # Debugging
            abort(404)

        return send_file(ppt_path, as_attachment=True)

    except Exception as e:
        print(f"Error: {e}")  # More detailed error logging
        return f"An error occurred while generating the PPT: {e}"

if __name__ == '__main__':
    app.run(debug=True)
